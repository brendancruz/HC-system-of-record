#!/usr/bin/env python3
"""Build the Diagnostics sector slide pack from the research-note content.

Generic template (no firm template loaded). Re-run after data refreshes:
    python3 sectors/build_diagnostics_deck.py
Output: sectors/diagnostics-deck.pptx

All figures mirror sectors/diagnostics-note.md and comps/diagnostics.csv
(as of 2026-06-13 prices/caps; FY2025 revenue from company releases).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Palette
NAVY = RGBColor(0x0F, 0x2A, 0x43)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
GREY = RGBColor(0x55, 0x5F, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB0, 0x30, 0x30)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def rect(s, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def set_run(r, text, size, color=NAVY, bold=False, italic=False):
    r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Calibri"


def header(s, eyebrow, title):
    rect(s, 0, 0, SW, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), SW, Inches(0.06), ACCENT)
    tb = box(s, Inches(0.55), Inches(0.16), Inches(12.2), Inches(0.95))
    tf = tb.text_frame
    p = tf.paragraphs[0]; set_run(p.add_run(), eyebrow.upper(), 11, ACCENT, bold=True)
    p2 = tf.add_paragraph(); set_run(p2.add_run(), title, 26, WHITE, bold=True)


def bullets(s, items, left, top, width, height, size=15, gap=10):
    tb = box(s, left, top, width, height)
    tf = tb.text_frame
    for i, (txt, lvl, *style) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.level = lvl
        bold = "b" in style
        color = RED if "r" in style else (ACCENT if "a" in style else NAVY)
        prefix = "" if lvl == 0 else "   "
        bullet = "" if "nodot" in style else ("▪  " if lvl == 0 else "–  ")
        set_run(p.add_run(), f"{prefix}{bullet}{txt}", size, color, bold=bold)


def footer(s, n):
    tb = box(s, Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    set_run(p.add_run(),
            "Diagnostics sector primer  |  BofA HC coverage prep  |  2026-06-14  "
            "|  Not investment advice  |  Sources: comps/diagnostics.csv, /companies/, /deals/",
            8, GREY)
    pn = box(s, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.35))
    pp = pn.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    set_run(pp.add_run(), str(n), 9, GREY)


# ---------------------------------------------------------------- Slide 1: title
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Inches(0.06), ACCENT)
tb = box(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.0))
tf = tb.text_frame
set_run(tf.paragraphs[0].add_run(), "HEALTHCARE DIAGNOSTICS", 14, ACCENT, bold=True)
p = tf.add_paragraph(); set_run(p.add_run(), "Sector Primer & Ideas Shortlist", 40, WHITE, bold=True)
sub = box(s, Inches(0.82), Inches(4.75), Inches(11.7), Inches(1.6))
sf = sub.text_frame
set_run(sf.paragraphs[0].add_run(),
        "Reimbursed, recurring molecular testing is re-rating around two TAMs - "
        "asymptomatic screening and MRD - while large-cap medtech sets the take-out floor at ~6-7x sales.",
        16, LIGHT)
p = sf.add_paragraph(); p.space_before = Pt(14)
set_run(p.add_run(), "BofA HC coverage prep   ·   2026-06-14   ·   7-name profiled universe", 12, RGBColor(0x9F,0xB3,0xC8))

# ---------------------------------------------------------------- Slide 2: exec summary
s = slide(); header(s, "Executive summary", "The valuation question is the whole story")
bullets(s, [
    ("Same income statement (high-growth, often pre-profit, reimbursement-dependent) earns "
     "1.6x to ~18x sales - depending entirely on which TAM story the market believes.", 0, "b"),
    ("Take-out anchor (~6-7x sales): Abbott bought Exact Sciences for ~$21B / $105/sh "
     "(closed 2026-03-23) at ~6.5x rev / ~52x adj. EBITDA - the price for profitable, recurring scale.", 0),
    ("Optionality premium (double-digit sales): Guardant (~18x) and Natera (~13x) trade on "
     "screening and MRD TAM, not current cash flow.", 0),
    ("Place any name between those two poles and you have the fluent read.", 0, "a", "b"),
    ("Profiled universe: ~$9.3B FY2025 revenue · ~$79.9B market cap · ~+30% rev-weighted growth "
     "· ~8.6x aggregate P/S (as of 2026-06-13; EXAS pinned at deal price).", 0),
    ("Recurring tell to watch: reimbursement + guideline news, NOT assay novelty.", 0, "r", "b"),
], Inches(0.55), Inches(1.5), Inches(12.2), Inches(5.3), size=16, gap=14)
footer(s, 2)

# ---------------------------------------------------------------- Slide 3: overview
s = slide(); header(s, "Industry overview", "How the money is made, and why growth is barbelled")
bullets(s, [
    ("Economics = volume x ASP. Moat = reimbursement coverage (a covered rate) + guideline "
     "inclusion (drives volume) + menu/scale (fixed-cost lab leverage). GM ~38% to ~74%.", 0, "b"),
    ("Growth is barbelled, not uniform:", 0, "a", "b"),
    ("Fast cohort: TEM +83% (Ambry-inflated; ~30% organic), WGS ~+40%, NTRA +36%, GH +31%.", 1),
    ("Slow cohort: EXAS +18% (scaled, profitable), FLGT +14%, NEO +10%.", 1),
    ("Four functional layers:", 0, "a", "b"),
    ("Screening (asymptomatic) - largest TAM, richest multiples (EXAS, GH/Shield).", 1),
    ("Liquid biopsy / therapy selection - the commercial core (Guardant360, Tempus, Natera).", 1),
    ("MRD / recurrence - highest-conviction growth after screening (Signatera, Reveal).", 1),
    ("Genomic & legacy lab - turnaround re-rate (WGS) to show-me value (NEO, FLGT).", 1),
], Inches(0.55), Inches(1.5), Inches(12.2), Inches(5.3), size=15, gap=9)
footer(s, 3)

# ---------------------------------------------------------------- Slide 4: landscape
s = slide(); header(s, "Competitive landscape", "Fought per battleground (use case + payor)")
# table
rows = [
    ("Battleground", "Leader(s)", "Challenger(s)"),
    ("CRC screening - stool", "Exact / Cologuard (Abbott)", "-"),
    ("CRC screening - blood", "Guardant / Shield", "Exact (Cancerguard), Grail"),
    ("Liquid biopsy - therapy sel.", "Guardant360, Foundation (Roche)", "Tempus, Natera"),
    ("MRD / recurrence", "Natera / Signatera", "Guardant Reveal, Tempus"),
    ("Hereditary / exome-genome", "GeneDx", "Natera, labs"),
    ("Broad / legacy molecular", "NeoGenomics, Fulgent", "Quest / Labcorp"),
]
t = s.shapes.add_table(len(rows), 3, Inches(0.55), Inches(1.45), Inches(8.4), Inches(3.6)).table
t.columns[0].width = Inches(2.9); t.columns[1].width = Inches(3.1); t.columns[2].width = Inches(2.4)
for ci, head in enumerate(rows[0]):
    c = t.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = NAVY
    p = c.text_frame.paragraphs[0]; set_run(p.add_run(), head, 11, WHITE, bold=True)
for ri in range(1, len(rows)):
    for ci in range(3):
        c = t.cell(ri, ci); c.fill.solid()
        c.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
        p = c.text_frame.paragraphs[0]
        set_run(p.add_run(), rows[ri][ci], 10, NAVY, bold=(ci == 0))
bullets(s, [
    ("Headline fight: blood-vs-stool CRC screening - Shield (GH) vs. Cologuard (Abbott); "
     "now a startup-vs-large-cap distribution mismatch.", 0, "b"),
    ("MRD is the cleanest secular winner after screening.", 0),
    ("Legacy tier is commoditizing toward reference-lab economics - more likely "
     "consolidated than re-rated.", 0),
    ("Bounding players (not profiled): Roche/Foundation, Grail, Quest/Labcorp, Abbott.", 0, "a"),
], Inches(9.15), Inches(1.5), Inches(3.7), Inches(5.0), size=12, gap=12)
footer(s, 4)

# ---------------------------------------------------------------- Slide 5: comps
s = slide(); header(s, "Peer comps spread", "P/S = mkt cap / FY2025 rev (EV/Sales proxy)")
comp = [
    ("Company", "Tkr", "Price", "Mkt Cap", "FY25 Rev", "Grw", "GM", "P/S", "26E P/S"),
    ("Natera", "NTRA", "$212.07", "~$30.4B", "$2.31B", "+36%", "~65%", "13.2x", "n/a"),
    ("Exact Sciences*", "EXAS", "$104.91", "~$20.0B", "$3.25B", "+18%", "~73%", "6.2x", "n/a"),
    ("Guardant Health", "GH", "$131.62", "~$17.45B", "~$0.97B", "+31%", "~65%", "17.9x", "13.3x"),
    ("Tempus AI", "TEM", "$47.82", "~$8.35B", "$1.27B", "+83%", "~52%", "6.6x", "5.3x"),
    ("GeneDx", "WGS", "$59.92", "~$1.78B", "~$0.43B", "~+40%", "~74%", "4.2x", "n/a"),
    ("NeoGenomics", "NEO", "$11.15", "~$1.43B", "$0.73B", "+10%", "~47%", "2.0x", "1.8x"),
    ("Fulgent Genetics", "FLGT", "$18.68", "~$0.53B", "$0.32B", "+14%", "~38%", "1.6x", "1.5x"),
]
t = s.shapes.add_table(len(comp), 9, Inches(0.55), Inches(1.45), Inches(12.2), Inches(3.5)).table
widths = [2.6, 0.9, 1.3, 1.5, 1.3, 1.1, 1.0, 1.0, 1.5]
for i, w in enumerate(widths):
    t.columns[i].width = Inches(w)
for ci, head in enumerate(comp[0]):
    c = t.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = NAVY
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
    set_run(p.add_run(), head, 11, WHITE, bold=True)
for ri in range(1, len(comp)):
    for ci in range(9):
        c = t.cell(ri, ci); c.fill.solid()
        c.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
        col = ACCENT if ci == 7 else NAVY
        set_run(p.add_run(), comp[ri][ci], 11, col, bold=(ci == 0 or ci == 7))
bullets(s, [
    ("* EXAS no longer freely traded (Abbott take-out) - P/S is the take-out benchmark. "
     "Summary P/S: median ~6.2x, mean ~7.4x, range 1.6x-17.9x.", 0),
    ("Watch FLGT: net cash exceeds market cap, so true EV/Sales is ~0 - 1.6x is NOT expensive.", 0, "r", "b"),
], Inches(0.55), Inches(5.15), Inches(12.2), Inches(1.7), size=13, gap=10)
footer(s, 5)

# ---------------------------------------------------------------- Slide 6: ideas
s = slide(); header(s, "Ideas shortlist", "Coverage thesis hooks (not recommendations)")
bullets(s, [
    ("Natera (NTRA) - MRD compounder.", 0, "a", "b"),
    ("Cleanest pure-play on MRD; Signatera driving +36% at $2.3B; 13.2x is growth-justified. "
     "Watch: MRD reimbursement breadth.", 1),
    ("Guardant (GH) - screening-optionality anchor.", 0, "a", "b"),
    ("~18x for Shield blood-based CRC optionality on a 25%+ oncology base. "
     "Watch: Shield coverage + guidelines; EBITDA negative.", 1),
    ("Tempus (TEM) - data re-rate wildcard.", 0, "a", "b"),
    ("+83% (~30% organic) yet only ~6.6x; data/Insights is the catalyst. Watch: organic growth + monetization.", 1),
    ("GeneDx (WGS) - convergence / turnaround.", 0, "a", "b"),
    ("Exome/genome +53-55% at ~74% adj. GM, only 4.2x; room to converge. Watch: sustained growth + FY25 actuals.", 1),
    ("Watch (not a hook): NeoGenomics / Fulgent - show-me / legacy; more likely consolidation targets.", 0, "b"),
], Inches(0.55), Inches(1.5), Inches(12.2), Inches(5.3), size=14, gap=8)
footer(s, 6)

# ---------------------------------------------------------------- Slide 7: how to deploy
s = slide(); header(s, "How to deploy in a coverage conversation", "The one-screen read")
bullets(s, [
    ("Open with the two-anchor frame: take-out ~6-7x (Abbott/EXAS) vs. screening/MRD premium "
     "(GH/NTRA) - then place each name on it in one sentence.", 0, "b"),
    ("Lead your differentiated takes with TEM (cheap-vs-growth gap, data optionality) and "
     "WGS (convergence candidate) - that is where you sound like you have done the work.", 0),
    ("The recurring tell is reimbursement + guidelines, not assay novelty.", 0, "a", "b"),
    ("Open items: refresh WGS to reported FY25 actuals; add net debt to get true EV/Sales; "
     "add EV/EBITDA for EBITDA-positive names; stand up Roche/Foundation, Grail, Quest/Labcorp profiles.", 0, "nodot"),
], Inches(0.55), Inches(1.7), Inches(12.2), Inches(4.5), size=16, gap=18)
footer(s, 7)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "diagnostics-deck.pptx")
prs.save(out)
print("wrote", out, "(", len(prs.slides._sldIdLst), "slides )")
